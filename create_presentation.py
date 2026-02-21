"""
Generate PowerPoint: TTS Systems for Voice Anti-Spoofing Research
Author: Tomas Acosta
Style: Clean white / gray / black with red, blue, green accents
Run: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────
# Colour palette  (white / gray / black primary; red, blue, green accents)
# ──────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF9, 0xFA, 0xFB)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY   = RGBColor(0x37, 0x41, 0x51)
BLACK       = RGBColor(0x11, 0x18, 0x27)
TITLE_BG    = RGBColor(0x1F, 0x29, 0x37)   # near-black header bar
CARD_BG     = RGBColor(0xF0, 0xF4, 0xF8)   # subtle blue-tinted card

BLUE        = RGBColor(0x1D, 0x4E, 0xD8)   # primary accent
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)   # badge background
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT= RGBColor(0xED, 0xE9, 0xFE)


# ──────────────────────────────────────────
# Core helpers
# ──────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, l, t, w, h, fill, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    return shp


def txt(slide, text, l, t, w, h,
        size=14, bold=False, italic=False,
        color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header_bar(slide, title):
    box(slide, 0, 0, 13.33, 1.05, TITLE_BG)
    box(slide, 0, 1.05, 13.33, 0.06, BLUE)   # blue underline
    txt(slide, title, 0.4, 0.1, 12.5, 0.85,
        size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def bullets(slide, items, l=0.55, t=1.3, w=12.3, base=14):
    """
    Items: list of strings.
      ## text  -> section subheader (blue, bold)
      >  text  -> right-column style (indented, gray)
      "  text  -> sub-bullet (indented gray)
      else     -> normal bullet (dark)
    """
    tb_box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(7.5 - t - 0.3))
    tf = tb_box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if item.startswith("##"):
            r = p.add_run()
            r.text = item[2:].strip()
            r.font.size  = Pt(base + 1)
            r.font.bold  = True
            r.font.color.rgb = BLUE
            p.space_before = Pt(8)
        elif item.startswith("  "):
            p.level = 1
            r = p.add_run()
            r.text = item.strip()
            r.font.size  = Pt(base - 1)
            r.font.color.rgb = DARK_GRAY
        elif item == "":
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(6)
        else:
            r = p.add_run()
            r.text = item
            r.font.size  = Pt(base)
            r.font.color.rgb = BLACK


def status_pill(slide, label, fill, text_color=WHITE, x=9.8, y=6.6, w=3.1, h=0.62):
    box(slide, x, y, w, h, fill)
    txt(slide, label, x, y + 0.08, w, h - 0.1,
        size=13, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def _cell(tbl, r, c, text, bg_fill, font_color=BLACK,
          size=11, bold=False, align=PP_ALIGN.CENTER):
    cell = tbl.cell(r, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_fill
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = font_color


# ──────────────────────────────────────────
# Slides
# ──────────────────────────────────────────
def s_title(prs):
    s = blank(prs)
    bg(s, BLACK)

    # Top red line
    box(s, 0, 0, 13.33, 0.12, RED)
    # Bottom white line
    box(s, 0, 7.38, 13.33, 0.12, WHITE)

    txt(s, "TTS SYSTEMS INVESTIGATION",
        0.6, 1.1, 12.1, 1.2, size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

    box(s, 2.5, 2.45, 8.33, 0.07, BLUE)

    txt(s, "Voice Anti-Spoofing Research · Latin American Spanish",
        0.6, 2.65, 12.1, 0.7, size=20, color=MID_GRAY,
        align=PP_ALIGN.CENTER)

    meta = ("Tomas Acosta\n"
            "Master's Thesis — Universidad de los Andes, Colombia\n"
            "February 2026  |  ml-server03: 4x NVIDIA A40 (46 GB VRAM each, CUDA 12.6)")
    txt(s, meta, 0.6, 3.55, 12.1, 1.6, size=14, color=MID_GRAY,
        align=PP_ALIGN.CENTER)

    txt(s, "Evaluating 6 state-of-the-art TTS systems for synthetic Spanish voice attack generation",
        0.6, 5.4, 12.1, 0.8, size=15, italic=True, color=MID_GRAY,
        align=PP_ALIGN.CENTER)


def s_exec_summary(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Executive Summary")

    bullets(s, [
        "## Purpose",
        "• Evaluate 6 modern TTS systems for generating synthetic Spanish voice attacks",
        "• Goal: create diverse training data for anti-spoofing detectors targeting Latin American speech",
        "",
        "## Three Evaluation Criteria",
        "• Spanish Language Quality  —  non-negotiable for Latin America scope",
        "• Implementation Feasibility  —  calendar timeline at 6 hrs/week, biweekly stakeholder meetings",
        "• Research Usefulness  —  attack sophistication, codec architecture diversity",
        "",
        "## Approach",
        "• No system is discarded for complexity alone — each is tiered by role and constraint",
        "• Complexity is a disadvantage that widens the timeline buffer, not a disqualifier",
        "• Hard restriction: English-only systems are repurposed for English baseline experiments",
        "",
        "## Key Finding",
        "• Fish Speech: Tier 1 (primary Spanish TTS)  |  Qwen3-TTS: Tier 2 (codec diversity)",
        "• Four additional systems assigned to Tiers 3–6 with specific roles and disclosed constraints",
    ], base=15)

    status_pill(s, "6 systems evaluated", BLUE, WHITE, x=9.5, y=6.6, w=3.5)


def s_criteria(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Evaluation Criteria")

    criteria = [
        (BLUE,  "1", "Spanish Language\nQuality",
         "MANDATORY\nLatin American phonetics,\naccent adaptability,\ntraining data volume"),
        (GREEN, "2", "Implementation\nFeasibility",
         "Calendar budget: 6 hrs/week\nBiweekly milestone meetings\nBuffer: min. 1.5 weeks\nper phase"),
        (RED,   "3", "Research\nUsefulness",
         "Attack sophistication\nCodec architecture diversity\nAcademic license compliance\nReproducibility"),
    ]
    for i, (color, num, title, desc) in enumerate(criteria):
        x = 0.4 + i * 4.3
        box(s, x, 1.3, 4.0, 5.8, LIGHT_GRAY, color, 1.5)
        box(s, x, 1.3, 4.0, 0.08, color)
        txt(s, num,   x + 0.2, 1.4,  0.6, 0.8, size=36, bold=True, color=color)
        txt(s, title, x + 0.2, 2.25, 3.6, 1.0, size=17, bold=True, color=BLACK, wrap=True)
        box(s, x + 0.2, 3.3, 3.6, 0.05, color)
        txt(s, desc,  x + 0.2, 3.45, 3.6, 3.4, size=13, color=DARK_GRAY, wrap=True)


def s_systems_overview(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Six Systems — Tiered by Role")

    systems = [
        ("Fish Speech\n(OpenAudio-S1)",   "4B params · RLHF",         "Tier 1 — Primary",          BLUE,   BLUE_LIGHT),
        ("Qwen3-TTS\n(Alibaba)",           "1.7B params · Dual-Track",  "Tier 2 — Secondary",        GREEN,  GREEN_LIGHT),
        ("CosyVoice 3.0\n(Alibaba)",       "1.5B params · CFM",         "Tier 6 — Deferred",         PURPLE, PURPLE_LIGHT),
        ("Chatterbox\n(Resemble.ai)",      "350M params · MIT",         "Tier 3 — Disclose Watermark",AMBER, AMBER_LIGHT),
        ("OuteTTS",                        "0.6B–1B · LLM-based",       "Tier 4 — Small Batch Only",  AMBER, AMBER_LIGHT),
        ("Nari Dia 1.6B\n(Nari Labs)",     "1.6B params · Dialogue",    "Tier 5 — English Baseline",  RED,   RED_LIGHT),
    ]

    xs = [0.3, 4.55, 8.8]
    ys = [1.35, 4.2]
    idx = 0
    for row_y in ys:
        for col_x in xs:
            name, detail, tier, color, bg_c = systems[idx]
            box(s, col_x, row_y, 3.9, 2.6, bg_c, color, 1.0)
            box(s, col_x, row_y, 3.9, 0.07, color)
            txt(s, name,   col_x + 0.2, row_y + 0.12, 3.5, 0.95,
                size=15, bold=True, color=BLACK, wrap=True)
            txt(s, detail, col_x + 0.2, row_y + 1.1,  3.5, 0.55,
                size=12, color=DARK_GRAY, wrap=True)
            box(s, col_x + 0.2, row_y + 1.75, 3.5, 0.6, color)
            txt(s, tier,   col_x + 0.2, row_y + 1.79, 3.5, 0.55,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
            idx += 1


def s_fish_overview(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Fish Speech (OpenAudio-S1) — Tier 1: Primary Spanish TTS")

    bullets(s, [
        "## What Is It?",
        "• 4 billion parameter open-source TTS system by Fish Audio",
        "• Trained with RLHF — optimises for human-perceived naturalness (not just technical accuracy)",
        "• Architecture: Dual Autoregressive (Dual-AR) + Firefly-GAN vocoder",
        "",
        "## Why It Ranks First",
        "• 20,000 hours of Spanish training data — same 'high-tier' as English and Chinese",
        "• Zero-shot voice cloning from 10–30 s reference audio (adapts to Latin American accents)",
        "• 24,906 GitHub stars · active maintenance · arXiv:2411.01156 · official Docker Compose",
        "• CC-BY-NC-SA-4.0 license explicitly permits academic / thesis research",
        "",
        "## Hardware on ml-server03",
        "• Requires 12 GB VRAM — we have 46 GB per A40 (26% utilisation per GPU)",
        "• Generation: ~2–3 s per 10-second sample  →  1,000 samples in ~33 minutes",
    ], base=15)

    status_pill(s, "Tier 1 — Primary", BLUE, WHITE, x=9.6, y=6.6, w=3.4)


def s_fish_tradeoffs(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Fish Speech — Advantages vs. Disadvantages")

    # Left column — advantages
    box(s, 0.35, 1.25, 6.0, 5.9, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 1.25, 6.0, 0.07, GREEN)
    txt(s, "Advantages", 0.55, 1.32, 5.7, 0.55,
        size=15, bold=True, color=GREEN)

    adv = [
        "Best validated Spanish quality in evaluation",
        "RLHF training: harder to detect than GAN-based TTS",
        "Cross-lingual voice cloning for LatAm accent adaptation",
        "Docker Compose: reproducible deployment",
        "Academic licence: no commercialisation risk",
        "Batch performance: 1,000 samples in ~33 minutes",
    ]
    y = 2.0
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.6, size=13, color=BLACK, wrap=True)
        y += 0.63

    # Right column — disadvantages
    box(s, 6.6, 1.25, 6.35, 5.9, RED_LIGHT, RED, 1.0)
    box(s, 6.6, 1.25, 6.35, 0.07, RED)
    txt(s, "Disadvantages / Caveats", 6.8, 1.32, 6.0, 0.55,
        size=15, bold=True, color=RED)

    dis = [
        "12 GB VRAM minimum (not a problem on A40, worth monitoring)",
        "No streaming: full audio generated before delivery (irrelevant for batch use)",
        "Voice cloning inconsistency on unusual accents — run validation tests first",
        "No specific LatAm sub-dialect benchmarks published",
        "Linux recommended — ml-server03 satisfies this",
    ]
    y = 2.0
    for d in dis:
        txt(s, "   " + d, 6.75, y, 6.0, 0.65, size=13, color=BLACK, wrap=True)
        y += 0.73

    status_pill(s, "Tier 1 — Primary", BLUE, WHITE, x=9.6, y=6.6, w=3.4)


def s_qwen(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Qwen3-TTS (Alibaba) — Tier 2: Codec Diversity")

    box(s, 0.35, 1.25, 6.0, 5.9, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 1.25, 6.0, 0.07, GREEN)
    txt(s, "Advantages", 0.55, 1.32, 5.7, 0.5, size=15, bold=True, color=GREEN)
    adv = [
        "Fastest inference: ~0.8 s per 10-second sample",
        "Apache 2.0 — unrestricted, no attribution required",
        "Easy install: pip install qwen-tts",
        "Different codec architecture from Fish Speech (diversity)",
        "FlashAttention 2 support: 30–40% speedup on A40",
        "1,000 samples in ~20–30 minutes (fastest of all systems)",
    ]
    y = 2.0
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.6, size=13, color=BLACK, wrap=True)
        y += 0.63

    box(s, 6.6, 1.25, 6.35, 5.9, AMBER_LIGHT, AMBER, 1.0)
    box(s, 6.6, 1.25, 6.35, 0.07, AMBER)
    txt(s, "Disadvantages / Restrictions", 6.8, 1.32, 6.0, 0.5, size=15, bold=True, color=AMBER)
    dis = [
        "Spanish is second-tier (paper: 'competitive' = not top-performing)",
        "Fine-tuning broken — cannot adapt to LatAm accents; base quality only",
        "Audio artifacts: truncated outputs, silent failures on long text",
        "Hard pin: transformers==4.57.3 — isolated conda environment required",
        "Dependency conflict with qwen-asr and similar ecosystem tools",
    ]
    y = 2.0
    for d in dis:
        txt(s, "   " + d, 6.75, y, 6.0, 0.65, size=13, color=BLACK, wrap=True)
        y += 0.73

    txt(s, "Role: Secondary (20% of samples).  Adds codec diversity; use Fish Speech for primary Spanish quality.",
        0.4, 6.55, 12.5, 0.7, size=12, italic=True, color=DARK_GRAY, wrap=True)
    status_pill(s, "Tier 2 — Secondary", GREEN, WHITE, x=9.6, y=6.6, w=3.4)


def s_cosyvoice(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "CosyVoice 3.0 (Alibaba) — Tier 6: Deferred (High Complexity)")

    box(s, 0.35, 1.25, 6.0, 5.9, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 1.25, 6.0, 0.07, GREEN)
    txt(s, "Advantages", 0.55, 1.32, 5.7, 0.5, size=15, bold=True, color=GREEN)
    adv = [
        "Largest training corpus: 1 million hours",
        "Apache 2.0 — unrestricted",
        "Highest published speaker similarity: 77.4%",
        "TensorRT-LLM: up to 4x speedup on A40",
        "Strong research backing (Alibaba Tongyi Lab)",
    ]
    y = 2.0
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.6, size=13, color=BLACK, wrap=True)
        y += 0.7

    box(s, 6.6, 1.25, 6.35, 5.9, AMBER_LIGHT, AMBER, 1.0)
    box(s, 6.6, 1.25, 6.35, 0.07, AMBER)
    txt(s, "Disadvantages / Restrictions", 6.8, 1.32, 6.0, 0.5, size=15, bold=True, color=AMBER)
    dis = [
        "vLLM version lock: only 0.9.0 OR 0.11.x+ — NOT 0.10.x",
          "  Users report environment corruption requiring fresh installs",
        "Deployment: 1–2 days (experienced engineer); buffer 2–3 weeks",
        "Zero Spanish benchmarks — quality is completely unvalidated",
        "Quality regression: community reports CosyVoice 2 sounds better",
        "Python 3.10 required; Matcha-TTS submodule dependency",
        "Go/no-go rule: if vLLM install >4 hours, defer indefinitely",
    ]
    y = 2.0
    for d in dis:
        if d.startswith("  "):
            txt(s, "     " + d.strip(), 6.75, y, 6.0, 0.55, size=12,
                color=DARK_GRAY, wrap=True)
            y += 0.55
        else:
            txt(s, "   " + d, 6.75, y, 6.0, 0.6, size=13, color=BLACK, wrap=True)
            y += 0.6

    status_pill(s, "Tier 6 — Deferred", PURPLE, WHITE, x=9.5, y=6.6, w=3.5)


def s_chatterbox(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Chatterbox (Resemble.ai) — Tier 3: Use with Disclosed Watermark")

    box(s, 0.35, 1.25, 6.0, 5.9, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 1.25, 6.0, 0.07, GREEN)
    txt(s, "Advantages", 0.55, 1.32, 5.7, 0.5, size=15, bold=True, color=GREEN)
    adv = [
        "Easiest setup: pip install chatterbox-tts (2–4 hours total)",
        "MIT license — truly open, no restrictions whatsoever",
        "Smallest VRAM: 8 GB (run 5+ instances in parallel on one A40)",
        "Outperformed ElevenLabs Turbo in blind tests (63.75% preference)",
        "Paralinguistic features: [laugh], [sigh], [hesitation] — adds realism",
        "Fastest path to first synthetic sample (prototyping)",
    ]
    y = 2.0
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.6, size=13, color=BLACK, wrap=True)
        y += 0.63

    box(s, 6.6, 1.25, 6.35, 5.9, RED_LIGHT, RED, 1.0)
    box(s, 6.6, 1.25, 6.35, 0.07, RED)
    txt(s, "Hard Restriction + Disadvantages", 6.8, 1.32, 6.0, 0.5,
        size=15, bold=True, color=RED)

    txt(s, "HARD RESTRICTION: Mandatory Perth watermark on all outputs.",
        6.75, 2.0, 6.0, 0.5, size=13, bold=True, color=RED, wrap=True)
    txt(s, ("Must be disclosed in thesis methodology.\n"
            "Use Chatterbox samples as isolated experimental group — "
            "detector must NOT mix watermarked + non-watermarked samples without flagging."),
        6.75, 2.52, 6.0, 1.1, size=12, color=DARK_GRAY, wrap=True)

    dis2 = [
        "Latency: claims <200 ms; delivers 300–600 ms (batch use: irrelevant)",
        "224 open issues / 34 commits — high bug density; test before scaling",
        "CPU inference broken despite being documented as supported",
        "No validated Latin American Spanish quality benchmarks",
    ]
    y = 3.75
    for d in dis2:
        txt(s, "   " + d, 6.75, y, 6.0, 0.6, size=12, color=BLACK, wrap=True)
        y += 0.62

    status_pill(s, "Tier 3 — Disclose Watermark", AMBER, WHITE, x=8.9, y=6.6, w=4.1)


def s_outetts(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "OuteTTS — Tier 4: Small-Scale Batch Only (Overnight Generation)")

    # Performance callout bar
    box(s, 0.35, 1.25, 12.6, 1.2, RED_LIGHT, RED, 1.5)
    txt(s, "Performance constraint: 3 min to generate 14 s on RTX 4090  |  1,000 samples = 1.5–2.5 days on A40",
        0.5, 1.35, 12.3, 0.6, size=15, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(s, "Solution: schedule as overnight / weekend server job — do not block any sprint on this.",
        0.5, 1.88, 12.3, 0.45, size=13, italic=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    box(s, 0.35, 2.65, 6.0, 4.5, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 2.65, 6.0, 0.07, GREEN)
    txt(s, "Advantages", 0.55, 2.72, 5.7, 0.5, size=15, bold=True, color=GREEN)
    adv = [
        "Spanish in 'high training data' tier (20–60 k hours)",
        "LLM-based architecture: maximum codec diversity",
        "Apache 2.0 license (0.6B variant)",
        "Speaker profiles stored as JSON — versionable",
        "Hosted API for small experiments ($0.0006/s)",
    ]
    y = 3.35
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.6, size=13, color=BLACK, wrap=True)
        y += 0.65

    box(s, 6.6, 2.65, 6.35, 4.5, AMBER_LIGHT, AMBER, 1.0)
    box(s, 6.6, 2.65, 6.35, 0.07, AMBER)
    txt(s, "Disadvantages / Restrictions", 6.8, 2.72, 6.0, 0.5, size=15, bold=True, color=AMBER)
    dis = [
        "76x slower than Fish Speech for batch generation",
        "100 samples/day API limit; $360 per 1,000 samples",
        "Sampling config fragility (64-token repetition window must be exact)",
        "Audio truncation mid-word; DAC codec sensitive to input quality",
        "Context limit: ~32 s effective generation per call",
    ]
    y = 3.35
    for d in dis:
        txt(s, "   " + d, 6.75, y, 6.0, 0.65, size=13, color=BLACK, wrap=True)
        y += 0.68

    status_pill(s, "Tier 4 — Small Batch Only", AMBER, WHITE, x=9.1, y=6.6, w=3.9)


def s_naridia(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Nari Dia 1.6B — Tier 5: English Baseline (Hard Restriction: No Spanish)")

    # Hard restriction callout
    box(s, 0.35, 1.25, 12.6, 1.0, RED_LIGHT, RED, 1.5)
    txt(s, "Hard Restriction: English language generation only.",
        0.55, 1.32, 12.3, 0.48, size=16, bold=True, color=RED)
    txt(s, "Cannot produce Spanish or Latin American Spanish. Repurposed for English baseline experiments.",
        0.55, 1.75, 12.3, 0.42, size=13, color=DARK_GRAY, wrap=True)

    box(s, 0.35, 2.45, 6.0, 4.7, GREEN_LIGHT, GREEN, 1.0)
    box(s, 0.35, 2.45, 6.0, 0.07, GREEN)
    txt(s, "Advantages (for English experiments)", 0.55, 2.52, 5.7, 0.55,
        size=14, bold=True, color=GREEN)
    adv = [
        "Ultra-realistic dialogue synthesis (ElevenLabs quality)",
        "Apache 2.0 — fully unrestricted",
        "Zero-shot cloning from 5 s reference audio",
        "Active development (April 2025, regularly updated)",
        "Paralinguistic: [laugh], [sigh], [cough] features",
        "Easy access: Hugging Face ZeroGPU (free demo)",
        "~10 GB VRAM — well within A40 capacity",
    ]
    y = 3.15
    for a in adv:
        txt(s, "   " + a, 0.5, y, 5.7, 0.58, size=13, color=BLACK, wrap=True)
        y += 0.58

    box(s, 6.6, 2.45, 6.35, 4.7, RED_LIGHT, RED, 1.0)
    box(s, 6.6, 2.45, 6.35, 0.07, RED)
    txt(s, "Hard Restriction + Disadvantages", 6.8, 2.52, 6.0, 0.55,
        size=14, bold=True, color=RED)

    txt(s, ("English only — no Spanish, no roadmap for Spanish.\n"
            "Next language targets: Asian languages (not Latin)."),
        6.75, 3.15, 6.0, 0.9, size=13, bold=True, color=RED, wrap=True)

    dis2 = [
        "Dialogue format ([S1]/[S2] tags) needs adaptation for single-speaker batch pipeline",
        "Optimised for 5–20 s dialogue chunks, not long-form reading",
        "Cannot contribute to core Latin American Spanish dataset",
    ]
    y = 4.15
    for d in dis2:
        txt(s, "   " + d, 6.75, y, 6.0, 0.7, size=13, color=BLACK, wrap=True)
        y += 0.72

    txt(s, "Thesis framing: 'Nari Dia used for English-language synthetic speech only — cross-lingual baseline.'",
        6.75, 5.85, 6.0, 0.85, size=11, italic=True, color=DARK_GRAY, wrap=True)

    status_pill(s, "Tier 5 — English Baseline", RED, WHITE, x=9.1, y=6.6, w=3.9)


def s_matrix(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Comparative Summary Matrix")

    rows, cols = 8, 7
    tbl = s.shapes.add_table(
        rows, cols,
        Inches(0.2), Inches(1.2),
        Inches(12.93), Inches(6.1)
    ).table

    headers = ["System", "Spanish Quality", "Work Hours", "Calendar\n(6 hrs/wk)",
               "VRAM", "License", "Status"]
    for ci, h in enumerate(headers):
        _cell(tbl, 0, ci, h, TITLE_BG, WHITE, size=12, bold=True)

    DATA = [
        ("Fish Speech",    "★★★★ Good\n(20 k hrs, validated)", "50–80 h",  "10–15 wks", "12 GB", "CC-BY-NC-SA-4.0\n(Academic OK)", "Tier 1 — Primary",   BLUE,   BLUE_LIGHT),
        ("Qwen3-TTS",      "★★★ Mediocre\n(not top-tier)",     "30–50 h",  "7–10 wks",  "4–8 GB","Apache 2.0",                    "Tier 2 — Secondary",  GREEN,  GREEN_LIGHT),
        ("Chatterbox",     "★★★ Unvalidated\n(23 languages)",  "20–35 h",  "5–8 wks",   "8 GB",  "MIT",                           "Tier 3 — Watermark\nDisclosed", AMBER, AMBER_LIGHT),
        ("OuteTTS",        "★★★ Adequate\n(20–60 k hrs)",      "15–30 h",  "4–7 wks\n+gen time","6–12 GB","Apache 2.0 (0.6B)",  "Tier 4 — Small\nBatch Only", AMBER, AMBER_LIGHT),
        ("Nari Dia 1.6B",  "English only\n(No Spanish)",       "15–25 h",  "4–6 wks",   "10 GB", "Apache 2.0",                    "Tier 5 — English\nBaseline",   RED,   RED_LIGHT),
        ("CosyVoice 3.0",  "★★ Unknown\n(no benchmarks)",      "40–80 h",  "10–18 wks\n(high variance)","8–16 GB","Apache 2.0", "Tier 6 — Deferred",   PURPLE, PURPLE_LIGHT),
    ]

    for ri, row in enumerate(DATA, start=1):
        *fields, label, color, bg_c = row
        row_bg = WHITE if ri % 2 == 0 else LIGHT_GRAY
        for ci, val in enumerate(fields):
            _cell(tbl, ri, ci, val, row_bg, BLACK, size=10,
                  align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
        _cell(tbl, ri, 6, label, bg_c, color, size=10, bold=True)


def s_performance(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Batch Generation Performance — 1,000 Samples (10 s each)")

    systems_perf = [
        ("Qwen3-TTS",       "20–30 min",    0.07,  BLUE),
        ("Fish Speech",     "~33 min",      0.09,  GREEN),
        ("CosyVoice 3.0",   "~35 min*",     0.10,  DARK_GRAY),
        ("Chatterbox",      "1.5–2 h",      0.30,  AMBER),
        ("OuteTTS (GPU)",   "1.5–2.5 DAYS", 1.0,   RED),
    ]

    bar_max_w = 8.5
    label_x   = 0.4
    bar_x     = 3.6

    for i, (name, time_str, ratio, color) in enumerate(systems_perf):
        y = 1.5 + i * 1.01
        box(s, label_x, y + 0.08, 3.0, 0.7, LIGHT_GRAY)
        txt(s, name, label_x + 0.1, y + 0.1, 2.8, 0.6, size=14, color=BLACK)
        bw = max(0.15, ratio * bar_max_w)
        box(s, bar_x, y + 0.1, bw, 0.6, color)
        txt(s, time_str, bar_x + bw + 0.2, y + 0.1, 2.5, 0.6,
            size=13, bold=True, color=color)

    txt(s, "* CosyVoice timing assumes vLLM environment installs cleanly (high uncertainty — verify before relying on this figure).",
        0.4, 6.78, 12.5, 0.5, size=11, italic=True, color=MID_GRAY, wrap=True)

    box(s, 0.4, 1.38, 12.9, 0.05, BLUE)


def s_spanish_ranking(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Spanish Language Support Ranking")

    rankings = [
        ("1st", "Fish Speech",   "★★★★ Good",    "20 k hrs · validated · high-tier alongside English & Chinese", BLUE),
        ("2nd", "Qwen3-TTS",     "★★★ Mediocre", "Paper: 'competitive' — not listed in top-performing group",    GREEN),
        ("3rd", "OuteTTS",       "★★★ Adequate", "20–60 k hrs, high-tier — practical constraints are limiting",  AMBER),
        ("4th", "Chatterbox",    "★★★ Unknown",  "23 languages listed — no benchmarks, no audio samples",        AMBER),
        ("5th", "CosyVoice 3.0", "★★ Unknown",   "No benchmarks at all — completely unvalidated for Spanish",    DARK_GRAY),
        ("---", "Nari Dia",      "English only", "Hard restriction — cannot produce Spanish in current release",  RED),
    ]

    for i, (medal, name, stars, evidence, color) in enumerate(rankings):
        y = 1.4 + i * 0.93
        box(s, 0.35, y, 12.6, 0.82, LIGHT_GRAY if i % 2 == 0 else WHITE, color, 0.8)
        box(s, 0.35, y, 0.08, 0.82, color)
        txt(s, medal, 0.5,  y + 0.15, 0.8, 0.55, size=13, bold=True, color=color)
        txt(s, name,  1.3,  y + 0.15, 2.5, 0.55, size=14, bold=True, color=BLACK)
        txt(s, stars, 3.85, y + 0.15, 2.4, 0.55, size=13, color=color)
        txt(s, evidence, 6.3, y + 0.15, 6.4, 0.55, size=12, color=DARK_GRAY, wrap=True)


def s_recommendation(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Recommended Integration Strategy")

    # Primary block
    box(s, 0.35, 1.25, 12.6, 1.9, BLUE_LIGHT, BLUE, 1.5)
    box(s, 0.35, 1.25, 0.1, 1.9, BLUE)
    txt(s, "Tier 1 — Fish Speech (Primary)", 0.6, 1.32, 8.0, 0.65,
        size=18, bold=True, color=BLUE)
    txt(s, ("Best validated Spanish quality (20 k training hours). RLHF-trained 4B-parameter synthesis. "
            "Anchor of the augmentation pipeline. 33 minutes for 1,000 samples on A40."),
        0.6, 1.93, 12.1, 0.95, size=13, color=BLACK, wrap=True)

    # Secondary block
    box(s, 0.35, 3.3, 12.6, 1.5, GREEN_LIGHT, GREEN, 1.5)
    box(s, 0.35, 3.3, 0.1, 1.5, GREEN)
    txt(s, "Tier 2 — Qwen3-TTS (Secondary, 20% of samples)", 0.6, 3.37, 10.0, 0.6,
        size=16, bold=True, color=GREEN)
    txt(s, ("Different codec architecture = improved detector generalisation. "
            "Fastest inference (0.8 s/sample). Contributes architectural diversity. "
            "Start only after Fish Speech is stable."),
        0.6, 3.95, 12.1, 0.75, size=13, color=BLACK, wrap=True)

    # Tiers 3–6 summary row
    tier_info = [
        ("Tier 3 — Chatterbox", "Disclose watermark,\nthen integrate", AMBER),
        ("Tier 4 — OuteTTS",    "Overnight batch\n(100–200 samples)", AMBER),
        ("Tier 5 — Nari Dia",   "English baseline\nexperiments only", RED),
        ("Tier 6 — CosyVoice",  "Defer until vLLM\ninstalls cleanly", PURPLE),
    ]
    xs = [0.35, 3.55, 6.75, 9.95]
    for i, (name, role, color) in enumerate(tier_info):
        x = xs[i]
        box(s, x, 5.0, 3.0, 1.85, LIGHT_GRAY, color, 0.8)
        box(s, x, 5.0, 3.0, 0.07, color)
        txt(s, name, x + 0.1, 5.08, 2.8, 0.6, size=12, bold=True, color=color, wrap=True)
        txt(s, role, x + 0.1, 5.7,  2.8, 0.95, size=12, color=DARK_GRAY, wrap=True)


def s_timeline(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Realistic Timeline — 6 hrs/week, Biweekly Stakeholder Meetings")

    # Working parameters
    box(s, 0.35, 1.25, 4.1, 2.3, LIGHT_GRAY)
    box(s, 0.35, 1.25, 4.1, 0.07, BLUE)
    txt(s, "Working Parameters", 0.5, 1.32, 3.9, 0.55, size=14, bold=True, color=BLUE)
    params = ["Hours/week: 6", "Meeting cadence: every 2 weeks",
              "Buffer (low complexity): 1.5 weeks", "Buffer (high complexity): 2–3 weeks",
              "Buffer (very high): 3+ weeks"]
    y = 2.0
    for p in params:
        txt(s, "  " + p, 0.45, y, 3.9, 0.42, size=12, color=BLACK)
        y += 0.42

    # Meeting milestone table
    tbl = s.shapes.add_table(
        8, 4,
        Inches(4.65), Inches(1.25),
        Inches(8.45), Inches(4.5)
    ).table

    for ci, h in enumerate(["Meeting", "Week", "Work Hours", "Deliverable"]):
        _cell(tbl, 0, ci, h, TITLE_BG, WHITE, size=12, bold=True)

    milestones = [
        ("M1", "2",  "6–10 h", "Fish Speech Docker running, 5 Spanish test samples"),
        ("M2", "4",  "6–10 h", "LatAm speaker profiles, quality validated"),
        ("M3", "6",  "6–8 h",  "FishSpeechAugmenter class + Qwen3 env ready"),
        ("M4", "8",  "6–8 h",  "End-to-end pipeline integrated"),
        ("M5", "10", "—",      "Buffer week (1.5-wk minimum)"),
        ("M6", "12", "6–8 h",  "1,000+ samples generated + validated"),
        ("M7", "14", "4–6 h",  "Thesis methodology section drafted"),
    ]

    for ri, (m, w, hrs, deliv) in enumerate(milestones, 1):
        bg_c = LIGHT_GRAY if ri % 2 == 0 else WHITE
        _cell(tbl, ri, 0, m,     bg_c, BLACK, size=11, bold=True)
        _cell(tbl, ri, 1, w,     bg_c, BLACK, size=11)
        _cell(tbl, ri, 2, hrs,   bg_c, BLACK, size=11)
        _cell(tbl, ri, 3, deliv, bg_c, BLACK, size=11, align=PP_ALIGN.LEFT)

    txt(s, "Tiers 1 & 2 combined: ~14 calendar weeks | Tiers 3–5 can overlap starting Week 17 | Tier 6 deferred",
        0.35, 5.88, 12.8, 0.5, size=12, italic=True, color=DARK_GRAY)


def s_why_tts(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Why TTS Augmentation Matters for Anti-Spoofing")

    box(s, 0.35, 1.25, 12.6, 1.4, BLUE_LIGHT, BLUE, 1.5)
    txt(s, ("Detectors trained only on traditional GAN-vocoder artifacts show 41.4% performance degradation\n"
            "against modern codec-based attacks — the dominant attack vector in 2025–2026."),
        0.55, 1.32, 12.2, 1.2, size=16, bold=True, color=BLUE, wrap=True)

    bullets(s, [
        "## Current Augmentation Pipeline (scientifically valid baseline)",
        "• 60% — Room Impulse Response + Background Noise",
        "• 30% — Codec Augmentation (MP3, AAC, Opus)",
        "• 10% — RawBoost signal perturbation",
        "",
        "## What TTS Adds",
        "• RLHF-trained synthesis (Fish Speech) is optimised for human perception — harder to detect",
        "• Diverse codec architectures (Dual-AR, Dual-Track, CFM) train a more generalisable detector",
        "• Cross-lingual experiments (Nari Dia in English) reveal language sensitivity of the detector",
        "",
        "• TTS is an enhancement that strengthens the research contribution",
          "  Existing pipeline remains the scientifically valid baseline if any TTS tier fails",
    ], t=2.85, base=15)


def s_professor_qa(prs):
    s = blank(prs)
    bg(s, WHITE)
    header_bar(s, "Anticipated Questions & Answers")

    qa = [
        ("Why Fish Speech and not CosyVoice with 1M training hours?",
         "CosyVoice has zero Spanish benchmarks — an unknown quality is not a better choice than a known good one. "
         "vLLM complexity is not a disqualifier but does demand a buffer; we defer CosyVoice to Tier 6 for that reason."),
        ("Why include Chatterbox if it watermarks all outputs?",
         "Mandatory watermarking is disclosed as a methodological restriction, not hidden. Chatterbox samples form "
         "an isolated experimental group; the thesis explicitly documents their watermark status."),
        ("How does this fit a 6-hour-per-week work schedule?",
         "Every tier has calendar timelines calculated at 6 hrs/week with biweekly meeting deliverables and minimum "
         "1.5-week buffers per phase. Tiers 1 + 2 span ~14 weeks; other tiers run in parallel or overlap."),
        ("What if Fish Speech Spanish quality is insufficient?",
         "Week 1-2 validation gate: 10–20 Spanish test samples evaluated before committing. "
         "If quality is inadequate, the existing pipeline (RIR + Codec + RawBoost) remains the primary baseline."),
    ]

    for i, (q, a) in enumerate(qa):
        y = 1.3 + i * 1.45
        box(s, 0.35, y, 12.6, 1.35, LIGHT_GRAY if i % 2 == 0 else OFF_WHITE, BLUE, 0.5)
        box(s, 0.35, y, 0.08, 1.35, BLUE)
        txt(s, "Q: " + q, 0.55, y + 0.07, 12.1, 0.55, size=13, bold=True, color=BLUE, wrap=True)
        txt(s, "A: " + a, 0.55, y + 0.62, 12.1, 0.65, size=12, color=BLACK, wrap=True)


def s_conclusion(prs):
    s = blank(prs)
    bg(s, BLACK)

    box(s, 0, 0,    13.33, 0.12, RED)
    box(s, 0, 7.38, 13.33, 0.12, WHITE)
    box(s, 0.35, 0.12, 0.06, 7.26, MID_GRAY)
    box(s, 12.92, 0.12, 0.06, 7.26, MID_GRAY)

    txt(s, "Conclusion", 0.6, 0.35, 12.1, 0.9,
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    box(s, 3.0, 1.35, 7.33, 0.06, BLUE)

    points = [
        ("Fish Speech is the primary TTS — best validated Spanish quality, stable codebase, 33 min/1k samples.", BLUE),
        ("Qwen3-TTS is the secondary contributor — codec diversity at fastest inference speed.", GREEN),
        ("Chatterbox adds rapid prototyping capacity — watermark disclosed as a methodological note.", AMBER),
        ("OuteTTS is viable at small scale — schedule generation as overnight server jobs.", AMBER),
        ("Nari Dia enables English cross-lingual baseline experiments — hard restriction documented.", RED),
        ("CosyVoice deferred — no Spanish benchmarks; pursue only if vLLM installs in under 4 hours.", MID_GRAY),
        ("All timelines calibrated to 6 hrs/week with minimum 1.5-week buffers and biweekly milestones.", WHITE),
    ]

    y = 1.58
    for pt, color in points:
        box(s, 0.55, y + 0.06, 0.18, 0.18, color)
        txt(s, pt, 0.85, y, 12.0, 0.5, size=13, color=WHITE, wrap=True)
        y += 0.68

    txt(s, "Tomas Acosta  |  Universidad de los Andes  |  February 2026",
        0.6, 6.85, 12.1, 0.42, size=12, color=MID_GRAY,
        align=PP_ALIGN.CENTER, italic=True)


# ──────────────────────────────────────────
# Build
# ──────────────────────────────────────────
def main():
    prs = new_prs()

    s_title(prs)               # 1
    s_exec_summary(prs)        # 2
    s_criteria(prs)            # 3
    s_systems_overview(prs)    # 4
    s_fish_overview(prs)       # 5
    s_fish_tradeoffs(prs)      # 6
    s_qwen(prs)                # 7
    s_cosyvoice(prs)           # 8
    s_chatterbox(prs)          # 9
    s_outetts(prs)             # 10
    s_naridia(prs)             # 11
    s_matrix(prs)              # 12
    s_performance(prs)         # 13
    s_spanish_ranking(prs)     # 14
    s_recommendation(prs)      # 15
    s_timeline(prs)            # 16
    s_why_tts(prs)             # 17
    s_professor_qa(prs)        # 18
    s_conclusion(prs)          # 19

    out = "TTS_Investigation_Presentation_v2.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
